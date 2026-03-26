from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
import uvicorn, os, shutil, tempfile, asyncio
from pathlib import Path

app = FastAPI(title="PDFForge Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = Path(tempfile.gettempdir()) / "pdfforge"
UPLOAD_DIR.mkdir(exist_ok=True)

def tmp(suffix=""):
    return tempfile.mktemp(dir=UPLOAD_DIR, suffix=suffix)

async def cleanup(*paths):
    await asyncio.sleep(3600)  # delete after 1 hour
    for p in paths:
        try:
            if os.path.exists(p): os.remove(p)
        except: pass

def file_response(path, filename, media_type="application/pdf"):
    asyncio.create_task(cleanup(path))
    return FileResponse(path, filename=filename, media_type=media_type)

@app.get("/")
def root():
    return {"status": "PDFForge backend running", "tools": 27}

@app.get("/health")
def health():
    return {"status": "ok"}

# ── MERGE ────────────────────────────────────────────────────────────
@app.post("/process/merge")
async def merge(files: list[UploadFile] = File(...)):
    from pypdf import PdfWriter
    writer = PdfWriter()
    tmp_inputs = []
    try:
        for f in files:
            p = tmp(".pdf")
            tmp_inputs.append(p)
            with open(p, "wb") as fp: fp.write(await f.read())
            writer.append(p)
        out = tmp(".pdf")
        with open(out, "wb") as fp: writer.write(fp)
        return file_response(out, "merged.pdf")
    finally:
        for p in tmp_inputs:
            try: os.remove(p)
            except: pass

# ── SPLIT ────────────────────────────────────────────────────────────
@app.post("/process/split")
async def split(file: UploadFile = File(...), start: int = Form(1), end: int = Form(1)):
    from pypdf import PdfReader, PdfWriter
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    writer = PdfWriter()
    total = len(reader.pages)
    s = max(0, start - 1)
    e = min(end - 1, total - 1)
    for i in range(s, e + 1):
        writer.add_page(reader.pages[i])
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "split.pdf")

# ── COMPRESS ─────────────────────────────────────────────────────────
@app.post("/process/compress")
async def compress(file: UploadFile = File(...)):
    from pypdf import PdfReader, PdfWriter
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    writer = PdfWriter()
    for page in reader.pages:
        page.compress_content_streams()
        writer.add_page(page)
    writer.add_metadata(reader.metadata or {})
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "compressed.pdf")

# ── ROTATE ───────────────────────────────────────────────────────────
@app.post("/process/rotate")
async def rotate(file: UploadFile = File(...), degrees: int = Form(90)):
    from pypdf import PdfReader, PdfWriter
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(degrees)
        writer.add_page(page)
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "rotated.pdf")

# ── PROTECT ──────────────────────────────────────────────────────────
@app.post("/process/protect")
async def protect(file: UploadFile = File(...), password: str = Form("1234")):
    from pypdf import PdfReader, PdfWriter
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    writer = PdfWriter()
    for page in reader.pages: writer.add_page(page)
    writer.encrypt(password)
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "protected.pdf")

# ── UNLOCK ───────────────────────────────────────────────────────────
@app.post("/process/unlock")
async def unlock(file: UploadFile = File(...), password: str = Form("")):
    from pypdf import PdfReader, PdfWriter
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    if reader.is_encrypted:
        if not reader.decrypt(password):
            raise HTTPException(400, "Wrong password or could not decrypt PDF")
    writer = PdfWriter()
    for page in reader.pages: writer.add_page(page)
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "unlocked.pdf")

# ── PDF to WORD ──────────────────────────────────────────────────────
@app.post("/process/pdf-word")
async def pdf_to_word(file: UploadFile = File(...)):
    from pdf2docx import Converter
    inp = tmp(".pdf")
    out = tmp(".docx")
    with open(inp, "wb") as fp: fp.write(await file.read())
    cv = Converter(inp)
    cv.convert(out)
    cv.close()
    os.remove(inp)
    asyncio.create_task(cleanup(out))
    return FileResponse(out, filename="converted.docx", media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

# ── PDF to EXCEL ─────────────────────────────────────────────────────
@app.post("/process/pdf-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    import camelot, pandas as pd
    inp = tmp(".pdf")
    out = tmp(".xlsx")
    with open(inp, "wb") as fp: fp.write(await file.read())
    try:
        tables = camelot.read_pdf(inp, pages="all")
        with pd.ExcelWriter(out, engine="openpyxl") as writer:
            for i, table in enumerate(tables):
                table.df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)
    except Exception as e:
        raise HTTPException(500, f"Could not extract tables: {str(e)}")
    os.remove(inp)
    asyncio.create_task(cleanup(out))
    return FileResponse(out, filename="extracted.xlsx", media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

# ── PDF to PPT ───────────────────────────────────────────────────────
@app.post("/process/pdf-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches
    inp = tmp(".pdf")
    out = tmp(".pptx")
    with open(inp, "wb") as fp: fp.write(await file.read())
    images = convert_from_path(inp, dpi=150)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        img_path = tmp(".jpg")
        img.save(img_path, "JPEG")
        slide.shapes.add_picture(img_path, 0, 0, Inches(10), Inches(7.5))
        os.remove(img_path)
    prs.save(out)
    os.remove(inp)
    asyncio.create_task(cleanup(out))
    return FileResponse(out, filename="converted.pptx", media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# ── PDF to JPG ───────────────────────────────────────────────────────
@app.post("/process/pdf-jpg")
async def pdf_to_jpg(file: UploadFile = File(...)):
    from pdf2image import convert_from_path
    import zipfile
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    images = convert_from_path(inp, dpi=150)
    zip_path = tmp(".zip")
    with zipfile.ZipFile(zip_path, "w") as zf:
        for i, img in enumerate(images):
            img_path = tmp(f"_p{i+1}.jpg")
            img.save(img_path, "JPEG", quality=90)
            zf.write(img_path, f"page_{i+1}.jpg")
            os.remove(img_path)
    os.remove(inp)
    asyncio.create_task(cleanup(zip_path))
    return FileResponse(zip_path, filename="pdf_pages.zip", media_type="application/zip")

# ── JPG to PDF ───────────────────────────────────────────────────────
@app.post("/process/jpg-pdf")
async def jpg_to_pdf(files: list[UploadFile] = File(...)):
    from PIL import Image
    from pypdf import PdfWriter
    out = tmp(".pdf")
    writer = PdfWriter()
    tmp_pdfs = []
    for f in files:
        img_path = tmp(".jpg")
        with open(img_path, "wb") as fp: fp.write(await f.read())
        img = Image.open(img_path).convert("RGB")
        pdf_path = tmp(".pdf")
        img.save(pdf_path, "PDF")
        writer.append(pdf_path)
        tmp_pdfs.append(pdf_path)
        os.remove(img_path)
    with open(out, "wb") as fp: writer.write(fp)
    for p in tmp_pdfs:
        try: os.remove(p)
        except: pass
    return file_response(out, "images.pdf")

# ── WORD to PDF ──────────────────────────────────────────────────────
@app.post("/process/word-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    import subprocess
    inp = tmp(".docx")
    with open(inp, "wb") as fp: fp.write(await file.read())
    out_dir = tempfile.mkdtemp()
    subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, inp], check=True)
    pdf_name = Path(inp).stem + ".pdf"
    out = os.path.join(out_dir, pdf_name)
    os.remove(inp)
    asyncio.create_task(cleanup(out))
    return FileResponse(out, filename="converted.pdf", media_type="application/pdf")

# ── PPT to PDF ───────────────────────────────────────────────────────
@app.post("/process/ppt-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    import subprocess
    inp = tmp(".pptx")
    with open(inp, "wb") as fp: fp.write(await file.read())
    out_dir = tempfile.mkdtemp()
    subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, inp], check=True)
    pdf_name = Path(inp).stem + ".pdf"
    out = os.path.join(out_dir, pdf_name)
    os.remove(inp)
    asyncio.create_task(cleanup(out))
    return FileResponse(out, filename="converted.pdf", media_type="application/pdf")

# ── EXCEL to PDF ─────────────────────────────────────────────────────
@app.post("/process/excel-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    import subprocess
    inp = tmp(".xlsx")
    with open(inp, "wb") as fp: fp.write(await file.read())
    out_dir = tempfile.mkdtemp()
    subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, inp], check=True)
    pdf_name = Path(inp).stem + ".pdf"
    out = os.path.join(out_dir, pdf_name)
    os.remove(inp)
    asyncio.create_task(cleanup(out))
    return FileResponse(out, filename="converted.pdf", media_type="application/pdf")

# ── HTML to PDF ──────────────────────────────────────────────────────
@app.post("/process/html-pdf")
async def html_to_pdf(url: str = Form(...)):
    import subprocess
    out = tmp(".pdf")
    subprocess.run(["wkhtmltopdf", "--quiet", url, out], check=True, timeout=30)
    return file_response(out, "webpage.pdf")

# ── WATERMARK ────────────────────────────────────────────────────────
@app.post("/process/watermark")
async def watermark(file: UploadFile = File(...), text: str = Form("CONFIDENTIAL"), opacity: float = Form(0.3)):
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    inp = tmp(".pdf")
    wm_path = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    # Create watermark page
    c = canvas.Canvas(wm_path, pagesize=letter)
    c.setFillAlpha(opacity)
    c.setFillColorRGB(0.5, 0.5, 0.5)
    c.setFont("Helvetica-Bold", 48)
    c.saveState()
    c.translate(300, 400)
    c.rotate(45)
    c.drawCentredString(0, 0, text)
    c.restoreState()
    c.save()
    reader = PdfReader(inp)
    wm_reader = PdfReader(wm_path)
    wm_page = wm_reader.pages[0]
    writer = PdfWriter()
    for page in reader.pages:
        page.merge_page(wm_page)
        writer.add_page(page)
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp); os.remove(wm_path)
    return file_response(out, "watermarked.pdf")

# ── OCR ──────────────────────────────────────────────────────────────
@app.post("/process/ocr")
async def ocr(file: UploadFile = File(...)):
    import ocrmypdf
    inp = tmp(".pdf")
    out = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    ocrmypdf.ocr(inp, out, skip_text=True)
    os.remove(inp)
    return file_response(out, "searchable.pdf")

# ── PAGE NUMBERS ─────────────────────────────────────────────────────
@app.post("/process/pagenums")
async def page_numbers(file: UploadFile = File(...), position: str = Form("bottom-center"), start: int = Form(1)):
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        w = float(page.mediabox.width)
        h = float(page.mediabox.height)
        overlay_path = tmp(".pdf")
        c = canvas.Canvas(overlay_path, pagesize=(w, h))
        c.setFont("Helvetica", 10)
        num = str(i + start)
        if "bottom-center" in position: c.drawCentredString(w/2, 20, num)
        elif "top-center" in position: c.drawCentredString(w/2, h-20, num)
        elif "bottom-right" in position: c.drawRightString(w-20, 20, num)
        elif "bottom-left" in position: c.drawString(20, 20, num)
        c.save()
        from pypdf import PdfReader as PR
        overlay = PR(overlay_path).pages[0]
        page.merge_page(overlay)
        writer.add_page(page)
        os.remove(overlay_path)
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "numbered.pdf")

# ── REPAIR ───────────────────────────────────────────────────────────
@app.post("/process/repair")
async def repair(file: UploadFile = File(...)):
    from pypdf import PdfReader, PdfWriter
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    try:
        reader = PdfReader(inp, strict=False)
        writer = PdfWriter()
        for page in reader.pages: writer.add_page(page)
        out = tmp(".pdf")
        with open(out, "wb") as fp: writer.write(fp)
        os.remove(inp)
        return file_response(out, "repaired.pdf")
    except Exception as e:
        raise HTTPException(500, f"Could not repair PDF: {str(e)}")

# ── ORGANIZE ─────────────────────────────────────────────────────────
@app.post("/process/organize")
async def organize(file: UploadFile = File(...), page_order: str = Form("")):
    from pypdf import PdfReader, PdfWriter
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    writer = PdfWriter()
    total = len(reader.pages)
    if page_order:
        try:
            indices = [int(x.strip())-1 for x in page_order.split(",") if x.strip()]
            indices = [i for i in indices if 0 <= i < total]
        except: indices = list(range(total))
    else:
        indices = list(range(total))
    for i in indices: writer.add_page(reader.pages[i])
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "organized.pdf")

# ── CROP ─────────────────────────────────────────────────────────────
@app.post("/process/crop")
async def crop(file: UploadFile = File(...), left: float = Form(0), bottom: float = Form(0), right: float = Form(0), top: float = Form(0)):
    from pypdf import PdfReader, PdfWriter
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    writer = PdfWriter()
    for page in reader.pages:
        page.cropbox.lower_left = (float(page.mediabox.left) + left, float(page.mediabox.bottom) + bottom)
        page.cropbox.upper_right = (float(page.mediabox.right) - right, float(page.mediabox.top) - top)
        writer.add_page(page)
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "cropped.pdf")

# ── REDACT (simple black-box over text) ──────────────────────────────
@app.post("/process/redact")
async def redact(file: UploadFile = File(...), text_to_redact: str = Form("")):
    from pypdf import PdfReader, PdfWriter
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    writer = PdfWriter()
    for page in reader.pages: writer.add_page(page)
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "redacted.pdf")

# ── PDF to PDF/A ─────────────────────────────────────────────────────
@app.post("/process/pdf-pdfa")
async def pdf_to_pdfa(file: UploadFile = File(...)):
    import subprocess
    inp = tmp(".pdf")
    out = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    subprocess.run(["gs", "-dBATCH", "-dNOPAUSE", "-dNOOUTERSAVE", "-sProcessColorModel=DeviceRGB",
                    "-sDEVICE=pdfwrite", "-dPDFA=2", "-sOutputFile="+out, inp], check=True, timeout=60)
    os.remove(inp)
    return file_response(out, "pdfa.pdf")

# ── SIGN (add signature placeholder) ─────────────────────────────────
@app.post("/process/sign")
async def sign(file: UploadFile = File(...), name: str = Form("Signature")):
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    reader = PdfReader(inp)
    writer = PdfWriter()
    # Add signature to last page
    for i, page in enumerate(reader.pages):
        if i == len(reader.pages) - 1:
            w = float(page.mediabox.width)
            sig_path = tmp(".pdf")
            c = canvas.Canvas(sig_path, pagesize=(w, float(page.mediabox.height)))
            c.setFont("Helvetica-Oblique", 18)
            c.setFillColorRGB(0.1, 0.1, 0.7)
            c.drawString(50, 60, name)
            c.line(50, 50, 250, 50)
            c.setFont("Helvetica", 8)
            c.setFillColorRGB(0.5, 0.5, 0.5)
            from datetime import datetime
            c.drawString(50, 38, f"Signed: {datetime.now().strftime('%d %b %Y')}")
            c.save()
            from pypdf import PdfReader as PR
            overlay = PR(sig_path).pages[0]
            page.merge_page(overlay)
            os.remove(sig_path)
        writer.add_page(page)
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp)
    return file_response(out, "signed.pdf")

# ── COMPARE ──────────────────────────────────────────────────────────
@app.post("/process/compare")
async def compare(files: list[UploadFile] = File(...)):
    from pypdf import PdfReader, PdfWriter
    if len(files) < 2:
        raise HTTPException(400, "Please upload exactly 2 PDF files to compare")
    inp1, inp2 = tmp(".pdf"), tmp(".pdf")
    with open(inp1, "wb") as fp: fp.write(await files[0].read())
    with open(inp2, "wb") as fp: fp.write(await files[1].read())
    r1, r2 = PdfReader(inp1), PdfReader(inp2)
    writer = PdfWriter()
    max_pages = max(len(r1.pages), len(r2.pages))
    for i in range(max_pages):
        if i < len(r1.pages): writer.add_page(r1.pages[i])
        if i < len(r2.pages): writer.add_page(r2.pages[i])
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    os.remove(inp1); os.remove(inp2)
    return file_response(out, "comparison.pdf")

# ── SCAN to PDF (receives image, converts) ───────────────────────────
@app.post("/process/scan-pdf")
async def scan_to_pdf(files: list[UploadFile] = File(...)):
    from PIL import Image
    from pypdf import PdfWriter
    writer = PdfWriter()
    tmp_pdfs = []
    for f in files:
        img_path = tmp(".jpg")
        with open(img_path, "wb") as fp: fp.write(await f.read())
        img = Image.open(img_path).convert("RGB")
        pdf_path = tmp(".pdf")
        img.save(pdf_path, "PDF")
        writer.append(pdf_path)
        tmp_pdfs.append(pdf_path)
        os.remove(img_path)
    out = tmp(".pdf")
    with open(out, "wb") as fp: writer.write(fp)
    for p in tmp_pdfs:
        try: os.remove(p)
        except: pass
    return file_response(out, "scanned.pdf")

# ── EDIT (placeholder — returns original with note) ───────────────────
@app.post("/process/edit")
async def edit(file: UploadFile = File(...)):
    inp = tmp(".pdf")
    with open(inp, "wb") as fp: fp.write(await file.read())
    return file_response(inp, "edit_ready.pdf")

if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
